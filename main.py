"""
日语词汇PPT图片增强工具
自动为PPT中的日语词汇搜索相关图片并添加到页面中
"""

import os
import requests  # pyright: ignore[reportMissingModuleSource]
from pptx import Presentation  # pyright: ignore[reportMissingImports]
from pptx.util import Inches, Pt  # pyright: ignore[reportMissingImports]
from pptx.enum.text import PP_ALIGN  # pyright: ignore[reportMissingImports]
from pptx.dml.color import RGBColor  # pyright: ignore[reportMissingImports]
import time
from urllib.parse import quote, urlencode
import json
import re
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    print("[WARN] Pillow未安装，无法转换WEBP格式图片。请运行: pip install Pillow")

class PPTImageEnhancer:
    def __init__(self, ppt_path, output_path=None, google_api_key=None, google_cse_id=None, 
                 google_ai_api_key=None, spark_api_key=None, spark_base_url=None, spark_model=None,
                 ollama_base_url=None, ollama_model=None, exa_api_key=None, serp_api_key=None, verbose=True):
        """
        初始化PPT图片增强器
        
        Args:
            ppt_path: 输入PPT文件路径
            output_path: 输出PPT文件路径（如果为None，则在原文件名后加_enhanced）
            google_api_key: Google Custom Search API Key（可选）
            google_cse_id: Google Custom Search Engine ID（可选）
            google_ai_api_key: Google AI (Gemini) API Key（可选，优先用于优化搜索关键词）
            spark_api_key: Spark AI API Key（可选，Gemini失败时使用）
            spark_base_url: Spark AI Base URL（可选）
            spark_model: Spark AI Model（可选）
            ollama_base_url: Ollama本地模型Base URL（可选，默认http://localhost:11434）
            ollama_model: Ollama模型名称（可选，默认llama3.2）
            exa_api_key: EXA API Key（可选）
            serp_api_key: Serp API Key（可选）
            verbose: 是否显示详细日志
        """
        self.ppt_path = ppt_path
        if output_path is None:
            base_name = os.path.splitext(ppt_path)[0]
            self.output_path = f"{base_name}_enhanced.pptx"
        else:
            self.output_path = output_path
        
        self.prs = Presentation(ppt_path)
        self.image_cache = {}  # 缓存已下载的图片
        self.google_api_key = google_api_key
        self.google_cse_id = google_cse_id
        self.google_ai_api_key = google_ai_api_key
        self.spark_api_key = spark_api_key
        self.spark_base_url = spark_base_url or "https://spark-api-open.xf-yun.com/v2"
        self.spark_model = spark_model or "spark-x"
        self.ollama_base_url = ollama_base_url or "http://localhost:11434"
        self.ollama_model = ollama_model or "llama3.2"
        self.exa_api_key = exa_api_key
        self.serp_api_key = serp_api_key
        self.verbose = verbose
        self.last_template_id = -1  # 记录上一页使用的模板ID，确保相邻页面不同
        self.optimized_keywords_cache = {}  # 缓存已优化的关键词，避免重复调用AI
        self.failed_urls = set()  # 记录失败的URL，避免重复尝试下载
        
        if self.verbose:
            if self.google_api_key and self.google_cse_id:
                print(f"[INFO] Google Custom Search API已配置")
            else:
                print(f"[INFO] Google API未配置")
            if self.exa_api_key:
                print(f"[INFO] EXA API已配置")
            if self.serp_api_key:
                print(f"[INFO] Serp API已配置")
            if self.google_ai_api_key:
                print(f"[INFO] Google AI (Gemini) API已配置，将优先用于优化搜索关键词")
            if self.ollama_base_url:
                print(f"[INFO] Ollama本地模型已配置: {self.ollama_base_url} (模型: {self.ollama_model})")
            if self.spark_api_key:
                print(f"[INFO] Spark AI API已配置，Gemini失败时将使用Spark")
            if not self.google_ai_api_key and not self.spark_api_key and not self.ollama_base_url:
                print(f"[INFO] AI API未配置，将直接使用原始关键词搜索")
        
    def search_images_google_api(self, keyword, count=2):
        """
        使用Google Custom Search API搜索图片
        
        Args:
            keyword: 搜索关键词
            count: 需要的图片数量
            
        Returns:
            图片URL列表
        """
        if not self.google_api_key or not self.google_cse_id:
            return []
        
        image_urls = []
        try:
            if self.verbose:
                print(f"    [DEBUG] 使用Google Custom Search API搜索: {keyword}")
            
            # Google Custom Search API
            url = "https://www.googleapis.com/customsearch/v1"
            params = {
                'key': self.google_api_key,
                'cx': self.google_cse_id,
                'q': keyword,
                'searchType': 'image',
                'num': min(count * 3, 10),  # 获取更多结果以便过滤WEBP
                'safe': 'active',
                'imgSize': 'xlarge',  # 使用xlarge获取更高分辨率图片
                'imgType': 'photo',
                'fileType': 'jpg,png'  # 限制为常见格式，避免webp
            }
            
            response = requests.get(url, params=params, timeout=15)
            
            if self.verbose:
                print(f"    [DEBUG] Google API响应状态: {response.status_code}")
            
            if response.status_code == 200:
                data = response.json()
                if 'items' in data:
                    for item in data['items']:
                        if len(image_urls) >= count:
                            break
                        img_url = item.get('link', '')
                        # 过滤WEBP格式（即使API设置了fileType，仍可能返回WEBP）
                        if img_url and not img_url.lower().endswith('.webp') and '.webp' not in img_url.lower():
                            image_urls.append(img_url)
                            if self.verbose:
                                print(f"    [DEBUG] 找到图片: {img_url[:60]}...")
                else:
                    if self.verbose:
                        print(f"    [DEBUG] Google API未返回结果")
            else:
                if self.verbose:
                    print(f"    [DEBUG] Google API错误: {response.status_code} - {response.text[:100]}")
                    
        except Exception as e:
            if self.verbose:
                print(f"    [DEBUG] Google API异常: {str(e)}")
        
        return image_urls
    
    def search_images_google_scrape(self, keyword, count=2):
        """
        通过爬取Google图片搜索结果获取图片（备用方案，不需要API key）
        
        Args:
            keyword: 搜索关键词
            count: 需要的图片数量
            
        Returns:
            图片URL列表
        """
        image_urls = []
        try:
            if self.verbose:
                print(f"    [DEBUG] 尝试爬取Google图片搜索结果: {keyword}")
            
            # 构建Google图片搜索URL
            search_url = f"https://www.google.com/search?tbm=isch&q={quote(keyword)}&safe=active"
            
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
                'Accept-Language': 'zh-CN,zh;q=0.9,en;q=0.8',
            }
            
            response = requests.get(search_url, headers=headers, timeout=15)
            
            if self.verbose:
                print(f"    [DEBUG] Google搜索结果响应状态: {response.status_code}")
            
            if response.status_code == 200:
                # 从HTML中提取图片URL
                html = response.text
                
                # 模式1：旧版Google图片结果中的"ou"原始图片URL
                # 优先提取原图URL（"ou"字段），避免使用缩略图
                pattern_ou = r'"ou":"([^"]+)"'  # 原始图片URL（高分辨率）
                matches_ou = re.findall(pattern_ou, html)
                
                if matches_ou:
                    for i, match in enumerate(matches_ou[:count]):
                        img_url = match.replace('\\u003d', '=').replace('\\/', '/')
                        # 跳过缩略图，优先使用原图
                        if 'encrypted-tbn0.gstatic.com' not in img_url or '=s' not in img_url:
                            image_urls.append(img_url)
                            if self.verbose:
                                print(f"    [DEBUG] 提取到原图 {len(image_urls)}: {img_url[:60]}...")
                        elif len(image_urls) < count:
                            # 如果是缩略图，尝试提取原图URL（去掉尺寸参数）
                            original_url = img_url.split('=s')[0] if '=s' in img_url else img_url
                            if original_url not in image_urls:
                                image_urls.append(original_url)
                                if self.verbose:
                                    print(f"    [DEBUG] 提取到原图（从缩略图转换） {len(image_urls)}: {original_url[:60]}...")
                
                # 如果原图不够，尝试其他模式
                if len(image_urls) < count:
                    if self.verbose:
                        print(f"    [DEBUG] 原图结果不足({len(image_urls)}/{count})，尝试其他模式")
                    
                    # 模式2：通用的 jpg/png/webp 链接（排除缩略图）
                    pattern2 = r'"(https://[^"]+\.(jpg|jpeg|png)[^"]*)"'
                    matches2 = re.findall(pattern2, html, re.IGNORECASE)
                    for url, ext in matches2:
                        if len(image_urls) >= count:
                            break
                        # 跳过缩略图
                        if 'encrypted-tbn0.gstatic.com' not in url or '=s' not in url:
                            if url not in image_urls:
                                image_urls.append(url)
                                if self.verbose:
                                    print(f"    [DEBUG] 提取到图片 {len(image_urls)} (备用): {url[:60]}...")
                    
                    # 模式3：从缩略图URL中提取原图（最后手段）
                    if len(image_urls) < count:
                        pattern3 = r'https://encrypted-tbn0\.gstatic\.com/images\?q=tbn:([^"&]+)'
                        matches3 = re.findall(pattern3, html)
                        for tbn_id in matches3[: (count - len(image_urls))]:
                            # 尝试构建可能的原图URL（但这不总是有效）
                            # 优先使用已经提取到的URL
                            pass
                    
                    # 模式4：兜底，从<img>标签中提取src（排除缩略图）
                    if len(image_urls) < count:
                        if self.verbose:
                            print(f"    [DEBUG] 前几种模式结果不足({len(image_urls)}/{count})，尝试从<img>标签中提取")
                        img_pattern = r'<img[^>]+src="(https://[^"]+)"'
                        img_matches = re.findall(img_pattern, html, re.IGNORECASE)
                        for url in img_matches:
                            if len(image_urls) >= count:
                                break
                            # 优先选择明显是原图的URL（包含常见图片域名）
                            if any(domain in url.lower() for domain in ['i.imgur.com', 'images.unsplash.com', 'pixabay.com', 'pexels.com']):
                                if url not in image_urls:
                                    image_urls.append(url)
                                    if self.verbose:
                                        print(f"    [DEBUG] 提取到图片 (img标签-原图): {url[:60]}...")
                            elif any(ext in url.lower() for ext in ['.jpg', '.jpeg', '.png']) and 'encrypted-tbn0' not in url:
                                if url not in image_urls:
                                    image_urls.append(url)
                                    if self.verbose:
                                        print(f"    [DEBUG] 提取到图片 (img标签): {url[:60]}...")
            else:
                if self.verbose:
                    print(f"    [DEBUG] Google搜索失败: {response.status_code}")
                    
        except Exception as e:
            if self.verbose:
                print(f"    [DEBUG] Google爬取异常: {str(e)}")
        
        return image_urls
    
    def optimize_search_keyword_with_gemini(self, japanese_text):
        """
        Use Google Gemini AI to optimize an image search keyword (优先使用).
        
        Args:
            japanese_text: Japanese text from the slide.
            
        Returns:
            Optimized search keyword (short English phrase), or None if optimization fails.
        """
        if not self.google_ai_api_key:
            return None
        
        try:
            prompt = (
                "You are helping to search images on Google Images.\n"
                "Given the following Japanese word or short phrase:\n"
                f"{japanese_text}\n\n"
                "Task:\n"
                "1. Generate ONE short English search query that will find images closely related\n"
                "   to the meaning of this Japanese word.\n"
                "2. The query should be at most 4–5 English words (very short).\n"
                "3. Output ONLY the English query text, without any explanation, quotes or extra words.\n"
                "4. If the Japanese word is a concrete thing (object, animal, food, place, action),\n"
                "   translate or describe it directly. If it is abstract, choose a concrete visual\n"
                "   concept that represents it (for example, a scene or object people can see).\n\n"
                "English search query:"
            )
            
            url = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash-exp:generateContent"
            params = {'key': self.google_ai_api_key}
            payload = {
                "contents": [{
                    "parts": [{"text": prompt}]
                }],
                "generationConfig": {
                    "temperature": 0.7,
                    "maxOutputTokens": 50
                }
            }
            
            if self.verbose:
                print(f"    [DEBUG] 调用Google Gemini AI优化关键词: {japanese_text}")
                print(f"    [DEBUG] Gemini API请求URL: {url}")
            
            response = requests.post(url, params=params, json=payload, timeout=30)
            response.raise_for_status()
            
            result = response.json()
            candidates = result.get("candidates", [])
            if not candidates:
                if self.verbose:
                    print(f"    [DEBUG] Gemini API未返回候选结果")
                return None
            
            content_parts = candidates[0].get("content", {}).get("parts", [])
            if not content_parts:
                if self.verbose:
                    print(f"    [DEBUG] Gemini API响应格式异常")
                return None
            
            content = content_parts[0].get("text", "").strip()
            
            if self.verbose:
                print(f"    [DEBUG] Gemini AI原始响应: {content[:100]}...")
            
            # Clean output
            optimized_keyword = content.split('\n')[0].strip()
            optimized_keyword = re.sub(r'^English search query[：:]\s*', '', optimized_keyword, flags=re.IGNORECASE)
            optimized_keyword = re.sub(r'^["\']|["\']$', '', optimized_keyword)
            
            if optimized_keyword and len(optimized_keyword) <= 40:
                if self.verbose:
                    print(f"    [DEBUG] Gemini AI优化后的关键词: {optimized_keyword}")
                return optimized_keyword[:40]
            else:
                if self.verbose:
                    print(f"    [DEBUG] Gemini AI返回的关键词格式不正确: {content}")
                return None
                
        except Exception as e:
            if self.verbose:
                print(f"    [DEBUG] Gemini AI优化关键词失败: {str(e)}")
            return None
    
    def optimize_search_keyword_with_spark(self, japanese_text):
        """
        Use Spark AI to optimize an image search keyword for Google Images.
        
        Args:
            japanese_text: Japanese text from the slide.
            
        Returns:
            Optimized search keyword (short English phrase, suitable for image search),
            or None if optimization fails.
        """
        if not self.spark_api_key:
            return None
        
        try:
            # Prompt in English, ask for a short English phrase for Google image search
            prompt = (
                "You are helping to search images on Google Images.\n"
                "Given the following Japanese word or short phrase:\n"
                f"{japanese_text}\n\n"
                "Task:\n"
                "1. Generate ONE short English search query that will find images closely related\n"
                "   to the meaning of this Japanese word.\n"
                "2. The query should be at most 4–5 English words (very short).\n"
                "3. Output ONLY the English query text, without any explanation, quotes or extra words.\n"
                "4. If the Japanese word is a concrete thing (object, animal, food, place, action),\n"
                "   translate or describe it directly. If it is abstract, choose a concrete visual\n"
                "   concept that represents it (for example, a scene or object people can see).\n\n"
                "English search query:"
            )
            
            headers = {
                "Authorization": f"Bearer {self.spark_api_key}",
                "Content-Type": "application/json",
            }
            
            base_url = self.spark_base_url.rstrip("/")
            endpoint = f"{base_url}/chat/completions"
            
            payload = {
                "model": self.spark_model,
                "messages": [{"role": "user", "content": prompt}],
                "stream": False,
                "temperature": 0.7,
                "max_tokens": 50,
            }
            
            if self.verbose:
                print(f"    [DEBUG] 调用Spark AI优化关键词: {japanese_text}")
                print(f"    [DEBUG] Spark API请求URL: {endpoint}")
            
            response = requests.post(endpoint, json=payload, headers=headers, timeout=30)
            response.raise_for_status()
            
            result = response.json()
            content = result.get("choices", [{}])[0].get("message", {}).get("content", "").strip()
            
            if self.verbose:
                print(f"    [DEBUG] Spark AI原始响应: {content[:100]}...")
            
            # Clean output: keep only the first line, strip quotes and extra labels
            optimized_keyword = content.split('\n')[0].strip()
            optimized_keyword = re.sub(r'^English search query[：:]\s*', '', optimized_keyword, flags=re.IGNORECASE)
            optimized_keyword = re.sub(r'^["\']|["\']$', '', optimized_keyword)
            
            # Limit length to avoid overly long prompts to Google
            if optimized_keyword and len(optimized_keyword) <= 40:
                if self.verbose:
                    print(f"    [DEBUG] Spark AI优化后的关键词: {optimized_keyword}")
                return optimized_keyword[:40]
            else:
                if self.verbose:
                    print(f"    [DEBUG] Spark AI返回的关键词格式不正确: {content}")
                return None
                
        except Exception as e:
            if self.verbose:
                print(f"    [DEBUG] Spark AI优化关键词失败: {str(e)}")
            return None
    
    def optimize_search_keyword_with_ollama(self, japanese_text):
        """
        Use Ollama local model to optimize an image search keyword.
        
        Args:
            japanese_text: Japanese text from the slide.
            
        Returns:
            Optimized search keyword (short English phrase), or None if optimization fails.
        """
        if not self.ollama_base_url:
            return None
        
        try:
            # 使用更简洁的prompt，减少处理时间
            prompt = (
                f"Translate this Japanese word to a short English image search query (2-4 words only): {japanese_text}\n"
                "Output only the English words, nothing else:"
            )
            
            base_url = self.ollama_base_url.rstrip("/")
            endpoint = f"{base_url}/api/generate"
            
            payload = {
                "model": self.ollama_model,
                "prompt": prompt,
                "stream": False,
                "options": {
                    "temperature": 0.3,  # 降低temperature，更确定性，更快
                    "num_predict": 20,   # 减少最大输出长度，只需要2-4个单词
                    "top_p": 0.9,        # 限制采样范围，加快速度
                    "top_k": 20          # 限制候选词数量
                }
            }
            
            if self.verbose:
                print(f"    [DEBUG] 调用Ollama本地模型优化关键词: {japanese_text}")
                print(f"    [DEBUG] Ollama API请求URL: {endpoint}")
            
            # 增加超时时间到120秒，因为首次调用需要加载模型
            response = requests.post(endpoint, json=payload, timeout=120)
            response.raise_for_status()
            
            result = response.json()
            content = result.get("response", "").strip()
            
            if self.verbose:
                print(f"    [DEBUG] Ollama原始响应: {content[:100]}...")
            
            # Clean output: keep only the first line, strip quotes and extra labels
            optimized_keyword = content.split('\n')[0].strip()
            optimized_keyword = re.sub(r'^English search query[：:]\s*', '', optimized_keyword, flags=re.IGNORECASE)
            optimized_keyword = re.sub(r'^["\']|["\']$', '', optimized_keyword)
            
            # Limit length to avoid overly long prompts
            if optimized_keyword and len(optimized_keyword) <= 40:
                if self.verbose:
                    print(f"    [DEBUG] Ollama优化后的关键词: {optimized_keyword}")
                return optimized_keyword[:40]
            else:
                if self.verbose:
                    print(f"    [DEBUG] Ollama返回的关键词格式不正确: {content}")
                return None
                
        except requests.exceptions.Timeout as e:
            if self.verbose:
                print(f"    [DEBUG] Ollama优化关键词超时: {str(e)}")
                print(f"    [DEBUG] 提示: 首次调用可能需要更长时间加载模型，跳过Ollama优化")
            # 超时不抛出异常，只返回None，让程序继续
            return None
        except requests.exceptions.ConnectionError as e:
            if self.verbose:
                print(f"    [DEBUG] Ollama连接错误: {str(e)}")
                print(f"    [DEBUG] 提示: Ollama服务可能未运行，跳过Ollama优化")
            return None
        except Exception as e:
            if self.verbose:
                print(f"    [DEBUG] Ollama优化关键词失败: {str(e)}")
            # 不抛出异常，只返回None，让程序继续
            return None
    
    def search_images_exa_api(self, keyword, count=2):
        """
        使用EXA API搜索图片
        
        Args:
            keyword: 搜索关键词
            count: 需要的图片数量
            
        Returns:
            图片URL列表
        """
        if not self.exa_api_key:
            return []
        
        image_urls = []
        try:
            if self.verbose:
                print(f"    [DEBUG] 使用EXA API搜索: {keyword}")
            
            # EXA API - 使用search端点
            url = "https://api.exa.ai/search"
            headers = {
                'x-api-key': self.exa_api_key,
                'Content-Type': 'application/json'
            }
            payload = {
                'query': keyword,
                'type': 'neural',
                'category': 'images',
                'num_results': min(count * 2, 10),  # 获取更多结果以便过滤
                'use_autoprompt': False,
                'contents': {'text': True, 'images': True}
            }
            
            response = requests.post(url, headers=headers, json=payload, timeout=15)
            
            if self.verbose:
                print(f"    [DEBUG] EXA API响应状态: {response.status_code}")
            
            if response.status_code == 200:
                data = response.json()
                if 'results' in data:
                    for item in data['results']:
                        if len(image_urls) >= count:
                            break
                        # EXA可能返回多种格式，尝试提取图片URL
                        img_url = item.get('url', '') or item.get('image_url', '') or item.get('image', '')
                        if img_url and not img_url.lower().endswith('.webp') and '.webp' not in img_url.lower():
                            image_urls.append(img_url)
                            if self.verbose:
                                print(f"    [DEBUG] 找到图片: {img_url[:60]}...")
                else:
                    if self.verbose:
                        print(f"    [DEBUG] EXA API未返回结果")
            else:
                if self.verbose:
                    print(f"    [DEBUG] EXA API错误: {response.status_code} - {response.text[:100]}")
                    
        except Exception as e:
            if self.verbose:
                print(f"    [DEBUG] EXA API异常: {str(e)}")
        
        return image_urls
    
    def search_images_serp_api(self, keyword, count=2):
        """
        使用Serp API搜索Google图片
        
        Args:
            keyword: 搜索关键词
            count: 需要的图片数量
            
        Returns:
            图片URL列表
        """
        if not self.serp_api_key:
            return []
        
        image_urls = []
        try:
            if self.verbose:
                print(f"    [DEBUG] 使用Serp API搜索: {keyword}")
            
            # Serp API for Google Images
            url = "https://serpapi.com/search"
            params = {
                'api_key': self.serp_api_key,
                'engine': 'google_images',
                'q': keyword,
                'num': min(count, 20),  # Serp API限制
                'safe': 'active',
                'ijn': 0  # 第一页
            }
            
            response = requests.get(url, params=params, timeout=15)
            
            if self.verbose:
                print(f"    [DEBUG] Serp API响应状态: {response.status_code}")
            
            if response.status_code == 200:
                data = response.json()
                if 'images_results' in data:
                    for item in data['images_results'][:count]:
                        img_url = item.get('original', '') or item.get('link', '')
                        if img_url and not img_url.lower().endswith('.webp') and '.webp' not in img_url.lower():
                            image_urls.append(img_url)
                            if self.verbose:
                                print(f"    [DEBUG] 找到图片: {img_url[:60]}...")
                else:
                    if self.verbose:
                        print(f"    [DEBUG] Serp API未返回结果")
            else:
                if self.verbose:
                    print(f"    [DEBUG] Serp API错误: {response.status_code} - {response.text[:100]}")
                    
        except Exception as e:
            if self.verbose:
                print(f"    [DEBUG] Serp API异常: {str(e)}")
        
        return image_urls
    
    def optimize_search_keyword_cached(self, keyword):
        """
        优化搜索关键词（带缓存，避免重复调用AI）
        
        Args:
            keyword: 原始关键词
            
        Returns:
            优化后的关键词，如果失败则返回None
        """
        # 检查缓存
        if keyword in self.optimized_keywords_cache:
            if self.verbose:
                print(f"    [DEBUG] 使用缓存的关键词优化结果: {self.optimized_keywords_cache[keyword]}")
            return self.optimized_keywords_cache[keyword]
        
        # 优先使用Gemini AI优化搜索关键词，失败则尝试Ollama，最后尝试Spark AI
        optimized_keyword = None
        if self.google_ai_api_key:
            optimized_keyword = self.optimize_search_keyword_with_gemini(keyword)
            if not optimized_keyword and self.ollama_base_url:
                if self.verbose:
                    print(f"    [DEBUG] Gemini优化失败，尝试使用Ollama本地模型")
                optimized_keyword = self.optimize_search_keyword_with_ollama(keyword)
            if not optimized_keyword and self.spark_api_key:
                if self.verbose:
                    print(f"    [DEBUG] Gemini和Ollama优化失败，尝试使用Spark AI")
                optimized_keyword = self.optimize_search_keyword_with_spark(keyword)
        elif self.ollama_base_url:
            optimized_keyword = self.optimize_search_keyword_with_ollama(keyword)
            if not optimized_keyword and self.spark_api_key:
                if self.verbose:
                    print(f"    [DEBUG] Ollama优化失败，尝试使用Spark AI")
                optimized_keyword = self.optimize_search_keyword_with_spark(keyword)
        elif self.spark_api_key:
            optimized_keyword = self.optimize_search_keyword_with_spark(keyword)
        
        # 缓存结果（包括None，避免重复尝试）
        self.optimized_keywords_cache[keyword] = optimized_keyword
        return optimized_keyword
    
    def search_images(self, keyword, count=2):
        """
        搜索与关键词相关的图片（优先使用Google图片搜索）
        
        Args:
            keyword: 搜索关键词（日语词汇）
            count: 需要的图片数量
            
        Returns:
            图片URL列表
        """
        image_urls = []
        
        # 清理关键词，提取主要词汇（去掉括号和假名）
        clean_keyword = keyword
        if '）' in clean_keyword or ')' in clean_keyword:
            parts = clean_keyword.split('）') if '）' in clean_keyword else clean_keyword.split(')')
            if len(parts) > 1:
                clean_keyword = parts[-1].strip()
        
        # 使用缓存的优化方法
        optimized_keyword = self.optimize_search_keyword_cached(clean_keyword)
        
        # 生成搜索关键词变体
        search_keywords = []
        if optimized_keyword:
            search_keywords.append(optimized_keyword)
        search_keywords.extend([
            clean_keyword,
            f"{clean_keyword} japanese",
            keyword,
        ])
        
        if self.verbose:
            print(f"    [DEBUG] 原始关键词: {clean_keyword}")
            if optimized_keyword:
                print(f"    [DEBUG] Spark AI优化后的关键词: {optimized_keyword}")
            print(f"    [DEBUG] 搜索关键词列表: {search_keywords}")
        
        # 方案1: 尝试使用Google Custom Search API
        if self.google_api_key and self.google_cse_id:
            for search_term in search_keywords:
                google_urls = self.search_images_google_api(search_term, count)
                # 过滤掉已知失败的URL
                google_urls = [url for url in google_urls if url not in self.failed_urls]
                image_urls.extend(google_urls)
                if len(image_urls) >= count:
                    break
        
        # 方案2: 如果Google API结果不足，尝试使用Serp API
        if len(image_urls) < count and self.serp_api_key:
            if self.verbose:
                print(f"    [DEBUG] Google API结果不足({len(image_urls)}/{count})，尝试使用Serp API")
            for search_term in search_keywords:
                serp_urls = self.search_images_serp_api(search_term, count - len(image_urls))
                # 过滤掉已知失败的URL
                serp_urls = [url for url in serp_urls if url not in self.failed_urls]
                image_urls.extend(serp_urls)
                if len(image_urls) >= count:
                    break
        
        # 方案3: 如果结果不足，尝试使用EXA API
        if len(image_urls) < count and self.exa_api_key:
            if self.verbose:
                print(f"    [DEBUG] API结果不足({len(image_urls)}/{count})，尝试使用EXA API")
            for search_term in search_keywords:
                exa_urls = self.search_images_exa_api(search_term, count - len(image_urls))
                # 过滤掉已知失败的URL
                exa_urls = [url for url in exa_urls if url not in self.failed_urls]
                image_urls.extend(exa_urls)
                if len(image_urls) >= count:
                    break
        
        # 方案4: 如果API结果不足，尝试爬取Google图片搜索结果
        if len(image_urls) < count:
            if self.verbose:
                print(f"    [DEBUG] API结果不足({len(image_urls)}/{count})，尝试爬取Google搜索结果")
            for search_term in search_keywords:
                scraped_urls = self.search_images_google_scrape(search_term, count - len(image_urls))
                # 过滤掉已知失败的URL
                scraped_urls = [url for url in scraped_urls if url not in self.failed_urls]
                image_urls.extend(scraped_urls)
                if len(image_urls) >= count:
                    break
        
        # 不再使用随机Picsum图片，只使用实际搜索结果
        if self.verbose:
            print(f"    [DEBUG] 总共找到 {len(image_urls)} 张图片（已过滤 {len(self.failed_urls)} 个已知失败的URL）")
        
        return image_urls[:count]
    
    def download_image(self, url, save_path, retry_count=3):
        """
        下载图片到本地（带重试机制和详细日志）
        
        Args:
            url: 图片URL
            save_path: 保存路径
            retry_count: 重试次数
            
        Returns:
            是否下载成功
        """
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
            'Accept': 'image/webp,image/apng,image/*,*/*;q=0.8',
            'Accept-Language': 'zh-CN,zh;q=0.9,en;q=0.8',
            'Referer': 'https://www.google.com/'
        }
        
        if self.verbose:
            print(f"    [DEBUG] 开始下载图片")
            print(f"    [DEBUG] URL: {url[:80]}...")
            print(f"    [DEBUG] 保存路径: {save_path}")
        
        # 为了避免在某个服务器上长时间反复重试，
        # 一旦出现超时或连接错误，立即放弃该URL，返回失败，
        # 由上层逻辑去尝试其他图片或重新搜索。
        for attempt in range(retry_count):
            try:
                if self.verbose and attempt > 0:
                    print(f"    [DEBUG] 重试 {attempt}/{retry_count}")
                
                start_time = time.time()
                response = requests.get(
                    url, 
                    timeout=30,  # 增加超时时间
                    allow_redirects=True, 
                    headers=headers,
                    stream=True  # 使用流式下载
                )
                elapsed_time = time.time() - start_time
                
                if self.verbose:
                    print(f"    [DEBUG] HTTP状态码: {response.status_code}")
                    print(f"    [DEBUG] 响应时间: {elapsed_time:.2f}秒")
                    print(f"    [DEBUG] Content-Type: {response.headers.get('Content-Type', '未知')}")
                    print(f"    [DEBUG] Content-Length: {response.headers.get('Content-Length', '未知')} bytes")
                
                if response.status_code == 200:
                    # 先根据 Content-Type 判断是否为WEBP，直接拒绝
                    content_type = response.headers.get('Content-Type', '').lower()
                    if 'image/webp' in content_type:
                        if self.verbose:
                            print(f"    [DEBUG] 检测到WEBP格式(Content-Type)，已跳过该图片: {url}")
                        self.failed_urls.add(url)  # 记录失败的URL
                        return False
                    
                    # 检查是否是有效的图片
                    content = response.content
                    content_size = len(content)
                    
                    if self.verbose:
                        print(f"    [DEBUG] 下载内容大小: {content_size} bytes")
                    
                    if content_size > 1000:  # 至少1KB
                        # 首先检查是否是HTML错误页面（常见错误响应）
                        content_str = content[:200].decode('utf-8', errors='ignore').lower()
                        if content_str.startswith('<!doctype') or content_str.startswith('<html') or '<html' in content_str[:100]:
                            if self.verbose:
                                print(f"    [DEBUG] ✗ 下载的内容是HTML页面而不是图片，已跳过: {url}")
                            self.failed_urls.add(url)  # 记录失败的URL
                            return False
                        
                        # 检查是否是图片格式
                        is_image = False
                        image_format = "未知"
                        
                        if content[:2] == b'\xff\xd8':
                            is_image = True
                            image_format = "JPEG"
                        elif content[:8] == b'\x89PNG\r\n\x1a\n':
                            is_image = True
                            image_format = "PNG"
                        elif content[:6] in [b'GIF87a', b'GIF89a']:
                            is_image = True
                            image_format = "GIF"
                        elif content[:12] == b'RIFF' and b'WEBP' in content[:20]:
                            # 明确识别为WEBP，直接跳过，不保存到本地
                            if self.verbose:
                                print(f"    [DEBUG] 检测到WEBP图片(内容特征)，已跳过该图片: {url}")
                            self.failed_urls.add(url)  # 记录失败的URL
                            return False
                        
                        if self.verbose:
                            print(f"    [DEBUG] 图片格式检测: {image_format}")
                        
                        # 只保存确认是图片格式的文件
                        if is_image:
                            with open(save_path, 'wb') as f:
                                f.write(content)
                            if self.verbose:
                                print(f"    [DEBUG] ✓ 图片下载成功: {save_path} ({content_size} bytes, {image_format})")
                            return True
                        else:
                            if self.verbose:
                                print(f"    [DEBUG] ✗ 内容不是有效的图片格式（不是JPEG/PNG/GIF）")
                            self.failed_urls.add(url)  # 记录失败的URL
                    else:
                        if self.verbose:
                            print(f"    [DEBUG] ✗ 内容太小 ({content_size} bytes < 1KB)")
                        self.failed_urls.add(url)  # 记录失败的URL
                else:
                    if self.verbose:
                        print(f"    [DEBUG] ✗ HTTP错误: {response.status_code}")
                    self.failed_urls.add(url)  # 记录失败的URL
                
            except requests.exceptions.Timeout as e:
                if self.verbose:
                    print(f"    [DEBUG] 超时异常: {str(e)}")
                    print(f"    [DEBUG] 立即放弃该图片URL，不再重试: {url}")
                self.failed_urls.add(url)  # 记录失败的URL
                # 直接放弃该URL，由上层逻辑尝试其他图片或重新搜索
                break
            except requests.exceptions.ConnectionError as e:
                if self.verbose:
                    print(f"    [DEBUG] 连接异常: {type(e).__name__}: {str(e)[:100]}")
                    print(f"    [DEBUG] 立即放弃该图片URL，不再重试: {url}")
                self.failed_urls.add(url)  # 记录失败的URL
                # 直接放弃该URL，由上层逻辑尝试其他图片或重新搜索
                break
            except Exception as e:
                if self.verbose:
                    print(f"    [DEBUG] 其他异常: {type(e).__name__}: {str(e)[:100]}")
                if attempt < retry_count - 1:
                    wait_time = 2 * (attempt + 1)
                    if self.verbose:
                        print(f"    [DEBUG] 等待 {wait_time} 秒后重试 ({attempt + 1}/{retry_count})...")
                    time.sleep(wait_time)
                    continue
                else:
                    if self.verbose:
                        print(f"    [DEBUG] ✗ 下载失败（异常，已重试{retry_count}次）")
                    self.failed_urls.add(url)  # 记录失败的URL
        
        if self.verbose:
            print(f"    [DEBUG] ✗ 所有重试均失败")
        self.failed_urls.add(url)  # 记录失败的URL
        return False
    
    def extract_text_from_slide(self, slide):
        """
        从幻灯片中提取所有文本
        
        Args:
            slide: 幻灯片对象
            
        Returns:
            文本列表
        """
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        return texts
    
    def print_progress(self, current, total):
        """
        在控制台打印处理进度条（带颜色）
        
        颜色规则（按百分比）：
        - < 20%: 红色
        - < 40%: 黄色
        - < 60%: 蓝色
        - < 80%: 青色
        - >= 80%: 绿色
        """
        if total <= 0:
            return
        
        percent = int(current * 100 / total)
        bar_length = 20  # 进度条长度（字符数）
        filled_length = int(bar_length * percent / 100)
        bar = "#" * filled_length + "-" * (bar_length - filled_length)
        
        # ANSI 颜色代码
        RED = "\033[31m"
        YELLOW = "\033[33m"
        BLUE = "\033[34m"
        CYAN = "\033[36m"
        GREEN = "\033[32m"
        RESET = "\033[0m"
        
        if percent < 20:
            color = RED
        elif percent < 40:
            color = YELLOW
        elif percent < 60:
            color = BLUE
        elif percent < 80:
            color = CYAN
        else:
            color = GREEN
        
        progress_text = f"[{bar}] {percent:3d}%  ({current}/{total} pages)"
        
        # 打印一行进度信息（不使用回车覆盖，避免与详细日志冲突）
        print(f"{color}{progress_text}{RESET}")
    
    def get_next_template_id(self):
        """
        获取下一个模板ID，确保与上一页不同
        
        Returns:
            模板ID (0-11)
        """
        # 定义12种不同的模板
        total_templates = 12
        
        # 生成候选模板列表（排除上一页使用的模板）
        available_templates = [i for i in range(total_templates) if i != self.last_template_id]
        
        # 随机选择（使用页面索引作为种子，确保可重复）
        import random
        # 使用当前处理的幻灯片数量作为种子，确保每次运行结果一致但相邻页面不同
        random.seed(len(self.prs.slides) * 7 + 13)  # 使用质数确保更好的随机性
        template_id = random.choice(available_templates)
        
        self.last_template_id = template_id
        return template_id
    
    def convert_image_format(self, image_path):
        """
        将图片转换为PowerPoint支持的格式（PNG或JPEG）
        
        Args:
            image_path: 原始图片路径
            
        Returns:
            转换后的图片路径（如果不需要转换则返回原路径）
        """
        # 首先检查文件是否存在
        if not os.path.exists(image_path):
            if self.verbose:
                print(f"    [DEBUG] 图片文件不存在: {image_path}")
            return image_path
        
        # 检查是否是文件路径（不是BytesIO或其他对象）
        if not isinstance(image_path, str):
            if self.verbose:
                print(f"    [DEBUG] 图片路径不是字符串类型: {type(image_path)}")
            return image_path
        
        if not PIL_AVAILABLE:
            # 如果没有PIL，无法安全转换WEBP等格式，直接返回原路径
            # 后续在 add_picture_safe 中会捕获不支持的格式错误，避免程序崩溃
            if self.verbose:
                print(f"    [DEBUG] Pillow未安装，无法转换图片格式: {image_path}")
            return image_path
        
        try:
            # 检查文件扩展名
            ext = os.path.splitext(image_path)[1].lower()
            supported_formats = ['.jpg', '.jpeg', '.png', '.bmp', '.gif', '.tiff', '.tif']
            
            if ext in supported_formats:
                # 已经是支持的格式，先验证文件确实是图片
                try:
                    # 尝试打开图片验证（verify会关闭文件，需要重新打开）
                    test_img = Image.open(image_path)
                    test_img.verify()
                    # verify后需要重新打开，因为verify会关闭文件
                    test_img = Image.open(image_path)
                    test_img.close()
                except Exception as e:
                    if self.verbose:
                        print(f"    [DEBUG] 文件扩展名是图片格式，但文件内容无效: {e}")
                        print(f"    [DEBUG] 错误类型: {type(e).__name__}")
                    # 如果验证失败，返回None而不是原路径，避免后续处理失败
                    return None
                return image_path
            
            # 需要转换格式
            if self.verbose:
                print(f"    [DEBUG] 转换图片格式: {ext} -> PNG")
            
            # 打开图片（这里会验证文件是否是有效的图片）
            # 确保 image_path 是字符串路径，不是 BytesIO 对象
            if not isinstance(image_path, str):
                raise ValueError(f"图片路径必须是字符串，但得到: {type(image_path)}")
            
            if not os.path.exists(image_path):
                raise FileNotFoundError(f"图片文件不存在: {image_path}")
            
            img = Image.open(image_path)
            
            # 如果是RGBA模式，转换为RGB（PNG支持透明度，但为了兼容性转为RGB）
            if img.mode == 'RGBA':
                # 创建白色背景
                background = Image.new('RGB', img.size, (255, 255, 255))
                background.paste(img, mask=img.split()[3])  # 使用alpha通道作为mask
                img = background
            elif img.mode not in ['RGB', 'L']:
                img = img.convert('RGB')
            
            # 生成新的文件名（PNG格式）
            new_path = os.path.splitext(image_path)[0] + '_converted.png'
            
            # 保存为PNG格式
            img.save(new_path, 'PNG', quality=95)
            
            if self.verbose:
                print(f"    [DEBUG] 图片已转换: {new_path}")
            
            return new_path
            
        except Exception as e:
            if self.verbose:
                print(f"    [DEBUG] 图片格式转换失败: {e}")
                print(f"    [DEBUG] 错误类型: {type(e).__name__}")
                print(f"    [DEBUG] 图片路径类型: {type(image_path).__name__}")
                if isinstance(image_path, str):
                    print(f"    [DEBUG] 图片路径: {image_path}")
                    print(f"    [DEBUG] 文件是否存在: {os.path.exists(image_path) if image_path else False}")
            # 转换失败，返回原路径（可能会失败，但至少尝试一下）
            # 但如果原路径也不是字符串，返回None让上层处理
            if isinstance(image_path, str):
                return image_path
            else:
                if self.verbose:
                    print(f"    [DEBUG] 图片路径不是字符串，无法返回: {type(image_path)}")
                return None
    
    def add_picture_safe(self, slide, image_path, x, y, max_width, max_height):
        """
        安全地添加图片到幻灯片，并在给定"框"中自适应缩放，保证不拉伸变形
        
        Args:
            slide: 幻灯片对象
            image_path: 图片路径
            x, y: 框的左上角位置
            max_width, max_height: 框的最大宽度和高度（图片在其中等比缩放、居中）
            
        Returns:
            添加的图片形状对象（或None，如果格式不支持）
        """
        # 首先验证输入
        if not isinstance(image_path, str):
            if self.verbose:
                print(f"    [DEBUG] 图片路径不是字符串类型: {type(image_path)}")
            return None
        
        if not os.path.exists(image_path):
            if self.verbose:
                print(f"    [DEBUG] 图片文件不存在: {image_path}")
            return None
        
        # 转换图片格式
        converted_path = self.convert_image_format(image_path)
        
        # 再次验证转换后的路径
        if converted_path is None:
            if self.verbose:
                print(f"    [DEBUG] 图片格式转换返回None，已跳过: {image_path}")
            return None
        
        if not isinstance(converted_path, str):
            if self.verbose:
                print(f"    [DEBUG] 转换后的图片路径不是字符串类型: {type(converted_path)}")
            return None
        
        if not os.path.exists(converted_path):
            if self.verbose:
                print(f"    [DEBUG] 转换后的图片文件不存在: {converted_path}")
            return None
        
        # 先以原始大小插入，再根据max_width/max_height做等比缩放和居中
        try:
            pic = slide.shapes.add_picture(converted_path, x, y)
        except (ValueError, IOError, OSError, FileNotFoundError) as e:
            msg = str(e)
            if "WEBP" in msg.upper():
                if self.verbose:
                    print(f"    [DEBUG] PowerPoint不支持WEBP图片，已跳过该图片: {converted_path}")
                    print("           如需使用WEBP图片，请安装 Pillow库 并重新运行：pip install Pillow")
                return None
            elif "cannot identify image file" in msg.lower() or "not a valid" in msg.lower() or "bytesio" in msg.lower():
                if self.verbose:
                    print(f"    [DEBUG] 无法识别图片文件（可能是无效的图片或HTML页面），已跳过: {converted_path}")
                    print(f"    [DEBUG] 错误详情: {msg}")
                return None
            # 其他错误继续抛出，便于排查
            if self.verbose:
                print(f"    [DEBUG] 添加图片时发生错误: {msg}")
                print(f"    [DEBUG] 错误类型: {type(e).__name__}")
            raise
        
        # 当前尺寸
        orig_w = pic.width
        orig_h = pic.height
        
        if orig_w <= 0 or orig_h <= 0:
            return pic
        
        # 计算等比缩放因子，使图片尽可能大地填充指定框（允许放大）
        scale_w = max_width / orig_w
        scale_h = max_height / orig_h
        scale = min(scale_w, scale_h)  # 允许放大，让图片尽可能大
        
        new_w = int(orig_w * scale)
        new_h = int(orig_h * scale)
        pic.width = new_w
        pic.height = new_h
        
        # 在框内居中对齐
        pic.left = int(x + (max_width - new_w) / 2)
        pic.top = int(y + (max_height - new_h) / 2)
        
        # 将图片置于底层（但仍然在文字遮罩之下）
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
        
        return pic
    
    def add_text_box(self, slide, text, x, y, width, height, bg_color=(255, 255, 255), transparency=0.15):
        """
        添加文字框的辅助函数（改进版：更好的可读性）
        
        Args:
            slide: 幻灯片对象
            text: 文字内容
            x, y, width, height: 位置和尺寸
            bg_color: 背景颜色RGB元组
            transparency: 透明度（降低默认透明度以提高可读性）
        """
        # 添加主文字框（白色背景，更不透明）
        text_box = slide.shapes.add_shape(1, x, y, width, height)
        fill = text_box.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*bg_color)
        fill.transparency = transparency  # 降低透明度，提高可读性
        
        # 添加边框（更粗的边框，提高对比度）
        line = text_box.line
        line.color.rgb = RGBColor(50, 50, 50)  # 深灰色边框
        line.width = Pt(4)  # 更粗的边框
        
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.4)
        text_frame.margin_right = Inches(0.4)
        text_frame.margin_top = Inches(0.3)
        text_frame.margin_bottom = Inches(0.3)
        
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(40)  # 增大字体
        p.font.bold = True
        # 使用深色文字，确保在任何背景下都清晰可见
        p.font.color.rgb = RGBColor(20, 20, 20)  # 深灰色，比纯黑更柔和但依然清晰
        p.alignment = PP_ALIGN.CENTER
        
        # 添加文字阴影效果（通过添加一个偏移的副本实现）
        # 注意：PowerPoint的python-pptx库不支持直接设置文字阴影，
        # 但我们可以通过调整文字颜色和背景来提高对比度
        
        return text_box
    
    def apply_template_0(self, slide, image_paths, original_texts):
        """模板0: 左右分割重叠（优化：文字框居中，更好的对比度）"""
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        combined_text = "\n".join(original_texts)
        
        if len(image_paths) > 0 and os.path.exists(image_paths[0]):
            self.add_picture_safe(slide, image_paths[0], Inches(0), Inches(0), slide_width * 0.6, slide_height)
        
        if len(image_paths) > 1 and os.path.exists(image_paths[1]):
            self.add_picture_safe(slide, image_paths[1], slide_width * 0.4, Inches(0), slide_width * 0.6, slide_height)
        
        # 文字框居中，降低透明度以提高可读性
        self.add_text_box(slide, combined_text, (slide_width - slide_width * 0.75) / 2, slide_height * 0.2, slide_width * 0.75, slide_height * 0.35, transparency=0.1)
    
    def apply_template_1(self, slide, image_paths, original_texts):
        """模板1: 上下分割（优化：文字框居中，更好的对比度）"""
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        combined_text = "\n".join(original_texts)
        
        if len(image_paths) > 0 and os.path.exists(image_paths[0]):
            self.add_picture_safe(slide, image_paths[0], Inches(0), Inches(0), slide_width, slide_height * 0.55)
        
        if len(image_paths) > 1 and os.path.exists(image_paths[1]):
            self.add_picture_safe(slide, image_paths[1], Inches(0), slide_height * 0.45, slide_width, slide_height * 0.55)
        
        # 文字框放在中间位置，降低透明度以提高可读性
        self.add_text_box(slide, combined_text, slide_width * 0.1, slide_height * 0.38, slide_width * 0.8, slide_height * 0.24, transparency=0.1)
    
    def apply_template_2(self, slide, image_paths, original_texts):
        """模板2: 对角线分割（优化：文字框位置和对比度）"""
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        combined_text = "\n".join(original_texts)
        
        if len(image_paths) > 0 and os.path.exists(image_paths[0]):
            self.add_picture_safe(slide, image_paths[0], Inches(0), Inches(0), slide_width * 0.7, slide_height * 0.7)
        
        if len(image_paths) > 1 and os.path.exists(image_paths[1]):
            self.add_picture_safe(slide, image_paths[1], slide_width * 0.3, slide_height * 0.3, slide_width * 0.7, slide_height * 0.7)
        
        # 文字框位置优化，降低透明度
        self.add_text_box(slide, combined_text, slide_width * 0.15, slide_height * 0.32, slide_width * 0.7, slide_height * 0.32, transparency=0.1)
    
    def apply_template_3(self, slide, image_paths, original_texts):
        """模板3: 中心聚焦（大图在中心，小图在角落）"""
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        combined_text = "\n".join(original_texts)
        
        if len(image_paths) > 0 and os.path.exists(image_paths[0]):
            self.add_picture_safe(slide, image_paths[0], slide_width * 0.15, slide_height * 0.15, slide_width * 0.7, slide_height * 0.7)
        
        if len(image_paths) > 1 and os.path.exists(image_paths[1]):
            self.add_picture_safe(slide, image_paths[1], slide_width * 0.7, slide_height * 0.05, slide_width * 0.25, slide_height * 0.25)
        
        self.add_text_box(slide, combined_text, slide_width * 0.2, slide_height * 0.75, slide_width * 0.6, slide_height * 0.2)
    
    def apply_template_4(self, slide, image_paths, original_texts):
        """模板4: 四象限布局"""
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        combined_text = "\n".join(original_texts)
        
        if len(image_paths) > 0 and os.path.exists(image_paths[0]):
            self.add_picture_safe(slide, image_paths[0], Inches(0), Inches(0), slide_width * 0.5, slide_height * 0.5)
        
        if len(image_paths) > 1 and os.path.exists(image_paths[1]):
            self.add_picture_safe(slide, image_paths[1], slide_width * 0.5, slide_height * 0.5, slide_width * 0.5, slide_height * 0.5)
        
        self.add_text_box(slide, combined_text, slide_width * 0.2, slide_height * 0.4, slide_width * 0.6, slide_height * 0.2)
    
    def apply_template_5(self, slide, image_paths, original_texts):
        """模板5: 左侧大图，右侧小图"""
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        combined_text = "\n".join(original_texts)
        
        if len(image_paths) > 0 and os.path.exists(image_paths[0]):
            self.add_picture_safe(slide, image_paths[0], Inches(0), Inches(0), slide_width * 0.65, slide_height)
        
        if len(image_paths) > 1 and os.path.exists(image_paths[1]):
            self.add_picture_safe(slide, image_paths[1], slide_width * 0.68, slide_height * 0.1, slide_width * 0.3, slide_height * 0.4)
        
        self.add_text_box(slide, combined_text, slide_width * 0.68, slide_height * 0.55, slide_width * 0.3, slide_height * 0.35)
    
    def apply_template_6(self, slide, image_paths, original_texts):
        """模板6: 顶部大图，底部小图"""
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        combined_text = "\n".join(original_texts)
        
        if len(image_paths) > 0 and os.path.exists(image_paths[0]):
            self.add_picture_safe(slide, image_paths[0], Inches(0), Inches(0), slide_width, slide_height * 0.65)
        
        if len(image_paths) > 1 and os.path.exists(image_paths[1]):
            self.add_picture_safe(slide, image_paths[1], slide_width * 0.1, slide_height * 0.68, slide_width * 0.4, slide_height * 0.3)
        
        self.add_text_box(slide, combined_text, slide_width * 0.55, slide_height * 0.68, slide_width * 0.4, slide_height * 0.3)
    
    def apply_template_7(self, slide, image_paths, original_texts):
        """模板7: 拼贴风格（两张图重叠交错）"""
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        combined_text = "\n".join(original_texts)
        
        if len(image_paths) > 0 and os.path.exists(image_paths[0]):
            self.add_picture_safe(slide, image_paths[0], slide_width * 0.1, slide_height * 0.1, slide_width * 0.5, slide_height * 0.6)
        
        if len(image_paths) > 1 and os.path.exists(image_paths[1]):
            self.add_picture_safe(slide, image_paths[1], slide_width * 0.4, slide_height * 0.3, slide_width * 0.5, slide_height * 0.6)
        
        self.add_text_box(slide, combined_text, slide_width * 0.15, slide_height * 0.72, slide_width * 0.7, slide_height * 0.25)
    
    def apply_template_8(self, slide, image_paths, original_texts):
        """模板8: 渐变叠加（两张图叠加，中间有遮罩，优化文字位置）"""
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        combined_text = "\n".join(original_texts)
        
        if len(image_paths) > 0 and os.path.exists(image_paths[0]):
            self.add_picture_safe(slide, image_paths[0], Inches(0), Inches(0), slide_width, slide_height)
        
        if len(image_paths) > 1 and os.path.exists(image_paths[1]):
            self.add_picture_safe(slide, image_paths[1], slide_width * 0.3, slide_height * 0.2, slide_width * 0.7, slide_height * 0.6)
            
        # 文字框放在底部，降低透明度以提高可读性
        self.add_text_box(slide, combined_text, slide_width * 0.1, slide_height * 0.72, slide_width * 0.8, slide_height * 0.24, transparency=0.1)
    
    def apply_template_9(self, slide, image_paths, original_texts):
        """模板9: 网格布局（2x2，但只用两张图）"""
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        combined_text = "\n".join(original_texts)
        
        if len(image_paths) > 0 and os.path.exists(image_paths[0]):
            self.add_picture_safe(slide, image_paths[0], slide_width * 0.05, slide_height * 0.05, slide_width * 0.45, slide_height * 0.45)
        
        if len(image_paths) > 1 and os.path.exists(image_paths[1]):
            self.add_picture_safe(slide, image_paths[1], slide_width * 0.5, slide_height * 0.5, slide_width * 0.45, slide_height * 0.45)
        
        self.add_text_box(slide, combined_text, slide_width * 0.2, slide_height * 0.7, slide_width * 0.6, slide_height * 0.25)
    
    def apply_template_10(self, slide, image_paths, original_texts):
        """模板10: 不规则分割（斜切效果）"""
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        combined_text = "\n".join(original_texts)
        
        if len(image_paths) > 0 and os.path.exists(image_paths[0]):
            self.add_picture_safe(slide, image_paths[0], Inches(0), Inches(0), slide_width * 0.55, slide_height)
        
        if len(image_paths) > 1 and os.path.exists(image_paths[1]):
            self.add_picture_safe(slide, image_paths[1], slide_width * 0.45, slide_height * 0.1, slide_width * 0.55, slide_height * 0.8)
        
        self.add_text_box(slide, combined_text, slide_width * 0.5, slide_height * 0.1, slide_width * 0.45, slide_height * 0.3)
    
    def apply_template_11(self, slide, image_paths, original_texts):
        """模板11: 画框风格（图片有边框，居中显示）"""
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        combined_text = "\n".join(original_texts)
        
        if len(image_paths) > 0 and os.path.exists(image_paths[0]):
            self.add_picture_safe(slide, image_paths[0], slide_width * 0.15, slide_height * 0.1, slide_width * 0.35, slide_height * 0.5)
            
            frame1 = slide.shapes.add_shape(1, slide_width * 0.15, slide_height * 0.1, slide_width * 0.35, slide_height * 0.5)
            frame1.fill.background()
            frame1.line.color.rgb = RGBColor(255, 255, 255)
            frame1.line.width = Pt(8)
            slide.shapes._spTree.remove(frame1._element)
            slide.shapes._spTree.insert(-1, frame1._element)
        
        if len(image_paths) > 1 and os.path.exists(image_paths[1]):
            self.add_picture_safe(slide, image_paths[1], slide_width * 0.5, slide_height * 0.1, slide_width * 0.35, slide_height * 0.5)
            
            frame2 = slide.shapes.add_shape(1, slide_width * 0.5, slide_height * 0.1, slide_width * 0.35, slide_height * 0.5)
            frame2.fill.background()
            frame2.line.color.rgb = RGBColor(255, 255, 255)
            frame2.line.width = Pt(8)
            slide.shapes._spTree.remove(frame2._element)
            slide.shapes._spTree.insert(-1, frame2._element)
        
        self.add_text_box(slide, combined_text, slide_width * 0.2, slide_height * 0.65, slide_width * 0.6, slide_height * 0.3)
    
    def add_creative_layout(self, slide, image_paths, original_texts, template_id=None):
        """
        为幻灯片添加创意布局的图片（12种不同模板）
        
        Args:
            slide: 幻灯片对象
            image_paths: 图片路径列表
            original_texts: 原始文本列表
            template_id: 模板ID（0-11），如果为None则自动选择
        """
        # 获取幻灯片尺寸
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        
        # 删除现有的所有内容（图片、形状、文本框等），避免重复显示
        shapes_to_remove = []
        for shape in slide.shapes:
            # 13=图片, 1=形状, 17=文本框, 14=占位符等
            if shape.shape_type in [13, 1, 17, 14]:
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
        
        # 如果没有指定模板ID，自动选择（确保与上一页不同）
        if template_id is None:
            template_id = self.get_next_template_id()
        
        if self.verbose:
            print(f"    [DEBUG] 使用模板 {template_id}/11")
        
        # 根据模板ID应用不同的布局
        template_methods = [
            self.apply_template_0, self.apply_template_1, self.apply_template_2,
            self.apply_template_3, self.apply_template_4, self.apply_template_5,
            self.apply_template_6, self.apply_template_7, self.apply_template_8,
            self.apply_template_9, self.apply_template_10, self.apply_template_11
        ]
        
        if 0 <= template_id < len(template_methods):
            template_methods[template_id](slide, image_paths, original_texts)
        else:
            # 如果模板ID无效，使用模板0
            self.apply_template_0(slide, image_paths, original_texts)
        
    def process_slides(self):
        """
        处理所有幻灯片
        """
        total_slides = len(self.prs.slides)
        print(f"开始处理PPT: {self.ppt_path}")
        print(f"共 {total_slides} 页幻灯片")
        
        # 创建临时图片目录
        temp_dir = "temp_images"
        os.makedirs(temp_dir, exist_ok=True)
        
        for idx, slide in enumerate(self.prs.slides):
            try:
                print(f"\n处理第 {idx + 1} 页...")
                
                # 提取文本
                texts = self.extract_text_from_slide(slide)
                if not texts:
                    print(f"  第 {idx + 1} 页没有找到文本，跳过")
                    self.print_progress(idx + 1, total_slides)
                    continue
                
                print(f"  提取的文本: {', '.join(texts)}")
                
                # 合并所有文本作为搜索关键词
                search_keyword = " ".join(texts)
                
                # 搜索图片
                print(f"  正在搜索图片...")
                image_urls = self.search_images(search_keyword, count=2)
                
                # 下载图片（跳过WEBP链接）
                image_paths = []
                for i, url in enumerate(image_urls):
                    image_path = os.path.join(temp_dir, f"slide_{idx}_img_{i}.jpg")
                    # 如果链接明显是WEBP，直接跳过，避免浪费请求
                    if url.lower().endswith('.webp'):
                        print(f"  下载图片 {i+1}/2... ✗ 跳过WEBP链接: {url}")
                        continue
                    print(f"  下载图片 {i+1}/2...", end='', flush=True)
                    if self.download_image(url, image_path):
                        image_paths.append(image_path)
                        print(f" ✓ 成功")
                    else:
                        print(f" ✗ 失败")
                
                # 强制要求必须有2张图片，如果不够则反复搜索直到找到2张
                max_retries = 3  # 最多重试3次（减少重试次数，避免超时）
                retry_count = 0
                
                while len(image_paths) < 2 and retry_count < max_retries:
                    if retry_count == 0 and not image_paths:
                        print(f"  ✗ 本轮图片下载失败，尝试使用首行文本重新搜索...")
                    elif len(image_paths) < 2:
                        print(f"  ✗ 图片数量不足({len(image_paths)}/2)，继续搜索...")
                    
                    # 使用第一行文本作为更精简的关键词，或者用AI优化关键词
                    retry_keyword = texts[0] if texts else search_keyword
                
                    # 如果配置了AI，尝试用AI优化关键词（使用缓存，避免重复调用）
                    # 只在第一次重试时尝试AI优化，避免API配额问题导致超时
                    if (self.google_ai_api_key or self.ollama_base_url or self.spark_api_key) and retry_count == 1:
                        if self.verbose:
                            print(f"    [DEBUG] 使用AI优化关键词（第{retry_count+1}次重试，使用缓存）")
                        # 使用缓存的优化方法，如果之前已经优化过，会直接返回缓存结果
                        optimized = self.optimize_search_keyword_cached(retry_keyword)
                        if optimized:
                            retry_keyword = optimized
                            if self.verbose:
                                print(f"    [DEBUG] AI优化后的关键词: {optimized}")
                    
                    # 重新搜索（search_images内部仍然优先用Google API和Google爬虫）
                    retry_urls = self.search_images(retry_keyword, count=2)
                    
                    # 再尝试下载（跳过已知失败的URL）
                    for i, url in enumerate(retry_urls):
                        if len(image_paths) >= 2:
                            break
                        image_path = os.path.join(temp_dir, f"slide_{idx}_retry{retry_count}_img_{i}.jpg")
                        # 如果链接明显是WEBP，直接跳过
                        if url.lower().endswith('.webp'):
                            if self.verbose:
                                print(f"  重新下载图片 {len(image_paths)+1}/2... ✗ 跳过WEBP链接: {url[:60]}...")
                            self.failed_urls.add(url)  # 记录失败的URL
                            continue
                        # 如果URL已经在失败列表中，直接跳过
                        if url in self.failed_urls:
                            if self.verbose:
                                print(f"  重新下载图片 {len(image_paths)+1}/2... ✗ 跳过已知失败的URL: {url[:60]}...")
                            continue
                        print(f"  重新下载图片 {len(image_paths)+1}/2...", end='', flush=True)
                        if self.download_image(url, image_path):
                            image_paths.append(image_path)
                            print(f" ✓ 成功")
                        else:
                            print(f" ✗ 失败")
                    
                    retry_count += 1
                
                # 如果图片下载成功，添加到幻灯片（必须使用2张，如果只有1张则重复使用）
                if image_paths:
                    # 确保有2张图片（如果只有1张，复制一份）
                    while len(image_paths) < 2:
                        import shutil
                        single_img = image_paths[0]
                        dup_path = single_img.replace('.jpg', '_dup.jpg')
                        shutil.copy2(single_img, dup_path)
                        image_paths.append(dup_path)
                        if self.verbose:
                            print(f"    [DEBUG] 图片不足2张，复制图片以补足: {dup_path}")
                    
                    image_paths = image_paths[:2]
                    print(f"  正在添加图片到幻灯片...")
                    self.add_creative_layout(slide, image_paths, texts)
                    print(f"  ✓ 第 {idx + 1} 页处理完成（模板ID: {self.last_template_id}，图片数: {len(image_paths)})")
                else:
                    print(f"  ✗ 第 {idx + 1} 页经过{max_retries}次搜索仍无法获得图片")
                    print(f"  ✗ 第 {idx + 1} 页所有搜索方案均失败，保留原始文字布局（不使用随机图片）")
                
                # 打印整体进度
                self.print_progress(idx + 1, total_slides)
                
            except Exception as e:
                # 单个页面处理失败，记录错误但继续处理其他页面
                error_msg = str(e)
                if self.verbose:
                    print(f"  ✗ 第 {idx + 1} 页处理失败: {error_msg}")
                    import traceback
                    print(f"  [DEBUG] 错误详情: {traceback.format_exc()}")
                else:
                    print(f"  ✗ 第 {idx + 1} 页处理失败，跳过")
                
                # 打印整体进度（即使失败也要更新）
                self.print_progress(idx + 1, total_slides)
        
        # 保存PPT
        print(f"\n保存处理后的PPT到: {self.output_path}")
        self.prs.save(self.output_path)
        print("✓ 处理完成！")
        
        # 清理临时文件夹
        import shutil
        try:
            if os.path.exists(temp_dir):
                shutil.rmtree(temp_dir)
                if self.verbose:
                    print(f"✓ 已清理临时文件夹: {temp_dir}")
        except Exception as e:
            if self.verbose:
                print(f"[WARN] 清理临时文件夹失败: {e}")
                print(f"      临时图片保存在 {temp_dir} 目录，可以手动删除")


def load_config():
    """
    从配置文件加载API配置
    
    Returns:
        (google_api_key, google_cse_id, google_ai_api_key, spark_api_key, spark_base_url, spark_model, 
         ollama_base_url, ollama_model, exa_api_key, serp_api_key) 元组
    """
    config_file = "config.json"
    google_api_key = None
    google_cse_id = None
    google_ai_api_key = None
    spark_api_key = None
    spark_base_url = None
    spark_model = None
    ollama_base_url = None
    ollama_model = None
    exa_api_key = None
    serp_api_key = None
    
    if os.path.exists(config_file):
        try:
            with open(config_file, 'r', encoding='utf-8') as f:
                config = json.load(f)
                google_api_key = config.get('google_api_key')
                google_cse_id = config.get('google_cse_id')
                google_ai_api_key = config.get('google_ai_api_key')  # Gemini API Key
                spark_api_key = config.get('spark_api_key')
                spark_base_url = config.get('spark_base_url', 'https://spark-api-open.xf-yun.com/v2')
                spark_model = config.get('spark_model', 'spark-x')
                ollama_base_url = config.get('ollama_base_url', 'http://localhost:11434')
                ollama_model = config.get('ollama_model', 'llama3.2')
                exa_api_key = config.get('exa_api_key')
                serp_api_key = config.get('serp_api_key')
                
                if google_api_key and google_cse_id:
                    print(f"[INFO] 已从配置文件加载Google Custom Search API设置")
                if google_ai_api_key:
                    print(f"[INFO] 已从配置文件加载Google AI (Gemini) API设置")
                if ollama_base_url:
                    print(f"[INFO] 已从配置文件加载Ollama本地模型设置: {ollama_base_url} (模型: {ollama_model})")
                if spark_api_key:
                    print(f"[INFO] 已从配置文件加载Spark AI API设置")
                if exa_api_key:
                    print(f"[INFO] 已从配置文件加载EXA API设置")
                if serp_api_key:
                    print(f"[INFO] 已从配置文件加载Serp API设置")
        except Exception as e:
            print(f"[WARN] 读取配置文件失败: {e}")
    
    return google_api_key, google_cse_id, google_ai_api_key, spark_api_key, spark_base_url, spark_model, ollama_base_url, ollama_model, exa_api_key, serp_api_key

def main():
    """
    主函数
    """
    print("=" * 50)
    print("日语词汇PPT图片增强工具")
    print("=" * 50)
    
    # 加载API配置
    google_api_key, google_cse_id, google_ai_api_key, spark_api_key, spark_base_url, spark_model, ollama_base_url, ollama_model, exa_api_key, serp_api_key = load_config()
    
    if not google_api_key or not google_cse_id:
        if not serp_api_key and not exa_api_key:
            print("\n[提示] Google图片搜索API、Serp API和EXA API均未配置")
            print("将尝试使用Google爬虫搜索图片（可能效果不佳）")
            print("建议配置Serp API、EXA API或Google Custom Search API")
        elif serp_api_key:
            print("\n[提示] Google图片搜索API未配置，将使用Serp API")
        elif exa_api_key:
            print("\n[提示] Google图片搜索API未配置，将使用EXA API")
    else:
        if serp_api_key or exa_api_key:
            print("\n[提示] 多个图片搜索API已配置")
        else:
            print("\n[提示] Google图片搜索API已配置")
    
    if not google_ai_api_key and not ollama_base_url and not spark_api_key:
        print("\n[提示] AI API未配置")
        print("将直接使用原始关键词搜索图片")
        print("如需使用AI优化搜索关键词，请在config.json中添加：")
        print("  - google_ai_api_key（优先，Gemini）")
        print("  - ollama_base_url（本地模型，推荐）")
        print("  - spark_api_key（备用）")
        print()
    elif google_ai_api_key:
        print("\n[提示] Google AI (Gemini) API已配置，将优先用于优化搜索关键词")
    elif ollama_base_url:
        print("\n[提示] Ollama本地模型已配置，将用于优化搜索关键词")
    elif spark_api_key:
        print("\n[提示] Spark AI API已配置，将用于优化搜索关键词")
    
    # 获取用户输入
    ppt_path = input("请输入PPT文件路径: ").strip().strip('"')
    
    if not os.path.exists(ppt_path):
        print(f"错误: 文件不存在: {ppt_path}")
        return
    
    if not ppt_path.endswith(('.pptx', '.ppt')):
        print("错误: 请提供.pptx或.ppt文件")
        return
    
    # 询问是否显示详细日志
    verbose_input = input("\n是否显示详细日志？(y/n，默认y): ").strip().lower()
    verbose = verbose_input != 'n'
    
    # 创建增强器并处理
    enhancer = PPTImageEnhancer(
        ppt_path, 
        google_api_key=google_api_key,
        google_cse_id=google_cse_id,
        google_ai_api_key=google_ai_api_key,
        spark_api_key=spark_api_key,
        spark_base_url=spark_base_url,
        spark_model=spark_model,
        ollama_base_url=ollama_base_url,
        ollama_model=ollama_model,
        exa_api_key=exa_api_key,
        serp_api_key=serp_api_key,
        verbose=verbose
    )
    enhancer.process_slides()


if __name__ == "__main__":
    main()

